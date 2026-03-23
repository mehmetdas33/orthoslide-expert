"""
PPTX Generation Engine
OrthoSlide Expert V2

Opens sunum.pptx template, replaces text placeholders with cephalometric
values, and inserts clinical images into specific slide placeholders.
"""
import os
import copy
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image

from ceph_logic import (
    REFERENCE_RANGES,
    EXCEL_ROW_MAP,
    evaluate_values,
    generate_diagnosis,
    get_placeholder_38_text,
    get_placeholder_40_text,
    get_placeholder_41_text,
    get_placeholder_80_text,
    get_placeholder_81_text,
    get_co_gn_ref_for_coa,
    get_ans_me_ref_for_coa,
    get_placeholder_70_text,
    get_placeholder_71_text,
    get_placeholder_95_text,
    get_placeholder_96_text,
    get_placeholder_97_text,
)


# ──────────────────────────────────────────────
# IMAGE SLOT → PPTX MAPPING
# Slides are 0-indexed.
# We map each UI slot to one or more [slide_index, placeholder_idx] targets.
# Slide 22 has exactly 9 placeholders (perfect for the 3x3 clinical composite).
# Slide 21 has exactly 4 placeholders (perfect for radiographs).
# ──────────────────────────────────────────────
IMAGE_SLOT_MAP = {
    # ── Slide index notes ─────────────────────────────────────────────────────
    # Template had slides 7 and 21 (1-based) removed by user.
    # 0-based shifts: orig 7→6, 8→7, 9→8, 10→9, 11→10, 12→11, 13→12,
    #                 14→13, 15→14, 16→15, 17→16, 18→17, orig 21→deleted, 22→20
    # ──────────────────────────────────────────────────────────────────────────

    # Kapak (Slayt 1, slide 0)
    "cover":             [{"slide": 0, "placeholder_idx": 2}],

    # Ekstraoral
    "frontal":           [{"slide": 2,  "placeholder_idx": 2},   # Slayt 3
                          {"slide": 20, "placeholder_idx": 10}], # Kompozit sol üst
    "frontal_smile":       [{"slide": 3,  "placeholder_idx": 2}],  # Slayt 4 (with midline)
    "frontal_smile_plain": [{"slide": 20, "placeholder_idx": 12}], # Kompozit sag üst (no midline)
    "profile":           [{"slide": 4,  "placeholder_idx": 2},   # Slayt 5 sol
                          {"slide": 20, "placeholder_idx": 16}], # Kompozit sol alt
    "profile_smile":     [{"slide": 4,  "placeholder_idx": 4}],  # Slayt 5 sag
    "three_quarter":     [{"slide": 20, "placeholder_idx": 18}], # Kompozit sag alt

    # Intraoral
    "intraoral_frontal": [{"slide": 5,  "placeholder_idx": 1},   # Slayt 6
                          {"slide": 20, "placeholder_idx": 14}], # Kompozit orta
    "intraoral_right":   [{"slide": 6,  "placeholder_idx": 1},   # Slayt 8 (orig 7→6)
                          {"slide": 20, "placeholder_idx": 15}], # Kompozit sol orta
    "intraoral_left":    [{"slide": 7,  "placeholder_idx": 1},   # Slayt 9 (orig 8→7)
                          {"slide": 20, "placeholder_idx": 13}], # Kompozit sag orta

    # Oklüzal — matched by top+left position (idx=4294967295)
    "upper_occlusal":    [{"slide": 9, "placeholder_idx": 4294967295,
                           "match_top": 0, "match_left": 4857752},   # Slayt 11 üst (orig 10→9)
                          {"slide": 20, "placeholder_idx": 11}],     # Kompozit üst orta
    "lower_occlusal":    [{"slide": 9, "placeholder_idx": 4294967295,
                           "match_top": 3429000, "match_left": 4857752},  # Slayt 11 alt
                          {"slide": 20, "placeholder_idx": 17}],          # Kompozit orta alt

    # Radyografi
    "panoramic":         [{"slide": 8,  "placeholder_idx": 1}],   # Slayt 10 (orig 9→8)
    "cephalometric":      [{"slide": 10, "placeholder_idx": 1, "fit": "contain"},   # Slayt 11
                          {"slide": 11, "placeholder_idx": 1, "fit": "contain"},    # Slayt 12
                          {"slide": 12, "placeholder_idx": 1, "fit": "contain"},    # Slayt 13 (uncropped)
                          {"slide": 15, "placeholder_idx": 1, "fit": "contain"}],   # Slayt 16
    "cephalometric_crop": [{"slide": 13, "placeholder_idx": 1, "fit": "contain"}],  # Slayt 14
    "cephalometric_line": [{"slide": 14, "placeholder_idx": 1, "fit": "contain"}],  # Slayt 15
    "ceph_raw":          [{"slide": 16, "placeholder_idx": 1, "fit": "contain"}],   # Slayt 18
    "wrist":             [{"slide": 17, "placeholder_idx": 2}],   # Slayt 19 (orig 18→17)
}


from PIL import Image, ImageOps

def insert_image_to_placeholder(slide, placeholder_idx, image_path, match_top=None, match_left=None, fit='cover'):
    """
    Insert an image into a placeholder on the given slide.
    fit='cover'   → center-crop to fill placeholder (default).
    fit='contain' → letterbox entire image with white background.
    When match_top / match_left are given, uses them as positional filters
    to disambiguate placeholders that share the same idx (e.g. idx=4294967295).
    """
    TOLERANCE = 300000  # ~8 mm in EMU
    placeholder = None
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        if shape.placeholder_format.idx != placeholder_idx:
            continue
        if match_top is not None and abs(shape.top - match_top) > TOLERANCE:
            continue
        if match_left is not None and abs(shape.left - match_left) > TOLERANCE:
            continue
        placeholder = shape
        break

    if placeholder is None:
        print(f"  Warning: Placeholder idx={placeholder_idx} not found on slide")
        return False

    left = placeholder.left
    top = placeholder.top
    width = placeholder.width
    height = placeholder.height

    import uuid
    with Image.open(image_path) as img:
        # Re-orient image according to EXIF data (fixes smartphone photos appearing sideways)
        img = ImageOps.exif_transpose(img)
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")

        img_w, img_h = img.size

        if fit == 'contain':
            # Contain: preserve full image, no crop, no white border.
            # Calculate display dimensions that fit within placeholder.
            img_ratio = img_w / img_h
            ph_ratio  = width / height  # EMU
            if img_ratio > ph_ratio:
                display_w = width
                display_h = int(width / img_ratio)
            else:
                display_h = height
                display_w = int(height * img_ratio)
            display_left = left + (width  - display_w) // 2
            display_top  = top  + (height - display_h) // 2

            # Resize image for quality (2000px max, content stays correct)
            scale = min(2000 / img_w, 2000 / img_h)
            resized = img.resize((max(1, int(img_w * scale)), max(1, int(img_h * scale))), Image.LANCZOS)
            temp_path = image_path + f"_tmp_{uuid.uuid4().hex[:6]}.jpg"
            resized.save(temp_path, "JPEG", quality=95)

            sp = placeholder._element
            sp.getparent().remove(sp)
            slide.shapes.add_picture(temp_path, display_left, display_top, display_w, display_h)
            try:
                os.remove(temp_path)
            except OSError:
                pass
            return True
        else:
            # cover: center-crop to fill placeholder
            img_ratio = img_w / img_h
            ph_ratio = width / height
            if img_ratio > ph_ratio:
                new_w = int(img_h * ph_ratio)
                left_crop = (img_w - new_w) // 2
                img = img.crop((left_crop, 0, left_crop + new_w, img_h))
            else:
                new_h = int(img_w / ph_ratio)
                top_crop = (img_h - new_h) // 2
                img = img.crop((0, top_crop, img_w, top_crop + new_h))

        temp_path = image_path + f"_tmp_{uuid.uuid4().hex[:6]}.jpg"
        img.save(temp_path, "JPEG", quality=95)

    sp = placeholder._element
    sp.getparent().remove(sp)

    slide.shapes.add_picture(
        temp_path,
        left, top,
        width, height
    )
    
    try:
        os.remove(temp_path)
    except OSError:
        pass

    return True


def replace_textbox_content(slide, shape_name, new_text):
    """
    Find a shape by name on the given slide and replace its text content.
    """
    for shape in slide.shapes:
        if shape.name == shape_name and shape.has_text_frame:
            tf = shape.text_frame
            if tf.paragraphs:
                para = tf.paragraphs[0]
                if para.runs:
                    for run in para.runs:
                        run.text = ""
                    para.runs[0].text = str(new_text)
                else:
                    para.text = str(new_text)
            return True
    return False


from pptx.dml.color import RGBColor

def find_and_replace_text(presentation, find_text, replace_text, status="normal", bold=False, align=None, font_size=None, font_name=None):
    """
    Search through all slides for text matching find_text and replace it.
    - status "high"/"low" applies red/blue colour.
    - bold=True makes the run bold.
    - font_size (int) sets the run font size in points.
    - font_name (str) sets the run font family.
    - align (PP_ALIGN constant) sets paragraph alignment on the matched paragraph.
    """
    color = None
    if status == "high":
        color = RGBColor(255, 0, 0)
    elif status == "low":
        color = RGBColor(0, 112, 192)

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if find_text in run.text:
                            run.text = run.text.replace(find_text, str(replace_text))
                            if color:
                                run.font.color.rgb = color
                            if bold:
                                run.font.bold = True
                            if font_size is not None:
                                run.font.size = Pt(font_size)
                            if font_name is not None:
                                run.font.name = font_name
                            if align is not None:
                                paragraph.alignment = align


def generate_pptx(
    template_path: str,
    output_path: str,
    ceph_data: dict,
    patient_info: dict = None,
    image_paths: dict = None,
    closing_video_path=None,
    pa_film_path=None,
):
    """
    Generate a completed PPTX from the template.
    """
    prs = Presentation(template_path)

    evaluated = evaluate_values(ceph_data)
    diagnosis = generate_diagnosis(ceph_data)

    placeholder_text_map = {}

    for i, key in enumerate(EXCEL_ROW_MAP):
        ph_num = i + 1
        val = ceph_data.get(key)
        
        status = "normal"
        for item in evaluated:
            if item["key"] == key:
                status = item["status"]
                break

        if val is not None:
            try:
                display_val = str(round(float(val)))
            except (ValueError, TypeError):
                display_val = str(val)
            placeholder_text_map[f"Placeholder {ph_num}"] = {
                "value": display_val, "status": status, "bold": True
            }

    # PH3: ANB = SNA − SNB (recalculated; overrides Excel row value)
    _sna_v = ceph_data.get("SNA")
    _snb_v = ceph_data.get("SNB")
    if _sna_v is not None and _snb_v is not None:
        try:
            _anb_calc = float(_sna_v) - float(_snb_v)
            _anb_ref  = REFERENCE_RANGES.get("ANB", {})
            _anb_s    = ("high" if _anb_calc > _anb_ref.get("high", 4)
                         else ("low" if _anb_calc < _anb_ref.get("low", 0) else "normal"))
            placeholder_text_map["Placeholder 3"] = {
                "value": str(round(_anb_calc)), "status": _anb_s, "bold": True
            }
        except (ValueError, TypeError):
            pass

    # PH38: ANB → Class I/II/III skeletal relationship (eski işlev)
    placeholder_text_map["Placeholder 38"] = {
        "value": get_placeholder_38_text(ceph_data), "status": "normal", "bold": True
    }
    # PH39: SNA+SNB → maxilla/mandible position combinations (Steiner prefix)
    _ph39_val = get_placeholder_80_text(ceph_data)
    placeholder_text_map["Placeholder 39"] = {
        "value": ("Steiner: " + _ph39_val) if _ph39_val else "", "status": "normal", "bold": True
    }
    placeholder_text_map["Placeholder 40"] = {
        "value": get_placeholder_40_text(ceph_data), "status": "normal", "bold": True
    }
    placeholder_text_map["Placeholder 41"] = {
        "value": get_placeholder_41_text(ceph_data), "status": "normal", "bold": True
    }
    # PH80: Wits+N-A → maxilla/mandible position combinations (McNamara prefix)
    _ph80_val = get_placeholder_81_text(ceph_data)
    placeholder_text_map["Placeholder 80"] = {
        "value": ("McNamara: " + _ph80_val) if _ph80_val else "", "status": "normal", "bold": True
    }
    # PH70/98: Upper incisor clinical text (U1-NA mm × U1-NA deg)
    _upper_inc = get_placeholder_70_text(ceph_data)
    placeholder_text_map["Placeholder 70"] = {"value": _upper_inc, "status": "normal", "bold": True}
    placeholder_text_map["Placeholder 98"] = {"value": _upper_inc, "status": "normal", "bold": True}

    # PH71/99: Lower incisor clinical text (L1-NB mm × L1-NB deg)
    _lower_inc = get_placeholder_71_text(ceph_data)
    placeholder_text_map["Placeholder 71"] = {"value": _lower_inc, "status": "normal", "bold": True}
    placeholder_text_map["Placeholder 99"] = {"value": _lower_inc, "status": "normal", "bold": True}

    # PH95: Nasolabial → increased / normal / decreased
    placeholder_text_map["Placeholder 95"] = {
        "value": get_placeholder_95_text(ceph_data), "status": "normal", "bold": True
    }
    # PH96: E-Upper lip → protruded / normal / retruded
    placeholder_text_map["Placeholder 96"] = {
        "value": get_placeholder_96_text(ceph_data), "status": "normal", "bold": True
    }
    # PH97: E-Lower lip → protruded / normal / retruded
    placeholder_text_map["Placeholder 97"] = {
        "value": get_placeholder_97_text(ceph_data), "status": "normal", "bold": True
    }

    # ── Jaw Length Analysis ──────────────────────────────────────────────────
    # PH85/86/87 = reference ranges (displayed on slide 14)
    # PH15/16/17 = actual measured values, colored by their reference ranges
    # PH90/91/92 = interpretation texts: increased / normal / decreased
    co_a_val = ceph_data.get("Co-A")
    co_a_float = None
    if co_a_val is not None:
        try:
            co_a_float = float(co_a_val)
        except (ValueError, TypeError):
            pass

    if co_a_float is not None:
        # PH86: Co-A reference = Co-A − N-A (dynamic, patient-specific)
        n_a_val = ceph_data.get("N-A")
        n_a_float = None
        if n_a_val is not None:
            try:
                n_a_float = float(n_a_val)
            except (ValueError, TypeError):
                pass

        if n_a_float is not None:
            co_a_reference = co_a_float - n_a_float
            placeholder_text_map["Placeholder 86"] = {
                "value": str(round(co_a_reference)), "status": "normal", "bold": True,
            }
            # PH16: Co-A actual value colored by PH86 reference
            _co_a_s = "high" if co_a_float > co_a_reference else ("low" if co_a_float < co_a_reference else "normal")
            placeholder_text_map["Placeholder 16"] = {
                "value": str(round(co_a_float)), "status": _co_a_s, "bold": True,
            }
        else:
            co_a_reference = None
            _co_a_s = "normal"

        # PH90: Co-A interpretation (increased/decreased/normal)
        placeholder_text_map["Placeholder 90"] = {
            "value": "increased" if _co_a_s == "high" else ("decreased" if _co_a_s == "low" else "normal"),
            "status": "normal", "bold": True,
        }

        # Dynamic lookup for Co-Gn and ANS-Me based on actual Co-A
        co_gn_low, co_gn_high = get_co_gn_ref_for_coa(co_a_float)
        ans_me_low, ans_me_high = get_ans_me_ref_for_coa(co_a_float)

        # PH87: Co-Gn reference range (dynamic)
        placeholder_text_map["Placeholder 87"] = {
            "value": f"{co_gn_low}-{co_gn_high}", "status": "normal", "bold": True,
        }
        # PH85: ANS-Me reference range (dynamic)
        placeholder_text_map["Placeholder 85"] = {
            "value": f"{ans_me_low}-{ans_me_high}", "status": "normal", "bold": True,
        }

        # PH17: Co-Gn actual value colored by dynamic range; PH91: interpretation
        co_gn_val = ceph_data.get("Co-Gn")
        if co_gn_val is not None:
            try:
                co_gn_f = float(co_gn_val)
                _co_gn_s = "high" if co_gn_f > co_gn_high else ("low" if co_gn_f < co_gn_low else "normal")
                placeholder_text_map["Placeholder 17"] = {
                    "value": str(round(co_gn_f)), "status": _co_gn_s, "bold": True,
                }
                placeholder_text_map["Placeholder 91"] = {
                    "value": "increased" if _co_gn_s == "high" else ("decreased" if _co_gn_s == "low" else "normal"),
                    "status": "normal", "bold": True,
                }
            except (ValueError, TypeError):
                pass

        # PH15: ANS-Me actual value colored by dynamic range; PH92: interpretation
        ans_me_raw = ceph_data.get("ANS-Me")
        if ans_me_raw is not None:
            try:
                ans_me_f = float(ans_me_raw)
                _ans_s = "high" if ans_me_f > ans_me_high else ("low" if ans_me_f < ans_me_low else "normal")
                placeholder_text_map["Placeholder 15"] = {
                    "value": str(round(ans_me_f)), "status": _ans_s, "bold": True,
                }
                placeholder_text_map["Placeholder 92"] = {
                    "value": "increased" if _ans_s == "high" else ("decreased" if _ans_s == "low" else "normal"),
                    "status": "normal", "bold": True,
                }
            except (ValueError, TypeError):
                pass

    # ── User-selected annotation answers (override any auto-generated values) ────
    # Arial Rounded MT Bold: 28pt for slides 1-10 and wrist (slayt 19),
    # 18pt for slide 11 (ph120-129).
    _ANNOTATION_FONT = "Arial Rounded MT Bold"
    # Slides 1-10 and slide 19: Arial Rounded MT Bold 28pt bold
    _ANNOTATION_28PT = [
        "ph101", "ph102", "ph103", "ph104", "ph105",   # class / profile
        "ph109",                                         # gülüş hattı (slayt 4)
        "ph130",                                         # wrist stage (slayt 19)
        "ph140", "ph141", "ph142", "ph143", "ph144",       # panoramik bulgular
    ]
    # Slide 11: Arial Rounded MT Bold 18pt bold
    _ANNOTATION_18PT = [
        "ph120", "ph121", "ph122",                       # maxilla counts (slayt 11)
        "ph123", "ph124", "ph125",                       # mandible counts (slayt 11)
        "ph126", "ph127", "ph128", "ph129",              # bolton (slayt 11)
    ]
    # Slide 6: plain (no bold, no font override — match template style)
    _ANNOTATION_PLAIN = [
        "ph110", "ph111",                                # midline texts (slayt 6)
        "ph112", "ph113",                                # overjet / overbite (slayt 6)
    ]
    if patient_info:
        # ── PH101 / PH103: molar & canine combined text ──────────────────────
        # Support both old format (separate 'I'/'II'/'III'/'No' per ph)
        # and new format (full sentence already built by frontend).
        _class_set = {'I', 'II', 'III', 'No'}
        _pi = dict(patient_info)  # mutable copy for transformation
        for _molar_ph, _canine_ph in [('ph101', 'ph102'), ('ph103', 'ph104')]:
            _mv = str(_pi.get(_molar_ph, '')).strip()
            _cv = str(_pi.get(_canine_ph, '')).strip()
            if _mv in _class_set or _cv in _class_set:
                # Old format: build combined sentence
                _m = _mv if _mv in _class_set else '?'
                _c = _cv if _cv in _class_set else '?'
                _pi[_molar_ph] = f"Class {_m} molar & Class {_c} canine relationship"
                _pi[_canine_ph] = ''   # canine ph is no longer needed separately
        patient_info = _pi

        for _ph_key in _ANNOTATION_28PT:
            _val = patient_info.get(_ph_key)
            if _val:
                _ph_num = _ph_key[2:]
                placeholder_text_map[f"Placeholder {_ph_num}"] = {
                    "value": str(_val), "status": "normal", "bold": True,
                    "font_size": 28, "font_name": _ANNOTATION_FONT,
                }
        for _ph_key in _ANNOTATION_18PT:
            _val = patient_info.get(_ph_key)
            if _val:
                _ph_num = _ph_key[2:]
                placeholder_text_map[f"Placeholder {_ph_num}"] = {
                    "value": str(_val), "status": "normal", "bold": True,
                    "font_size": 18, "font_name": _ANNOTATION_FONT,
                }
        for _ph_key in _ANNOTATION_PLAIN:
            _val = patient_info.get(_ph_key)
            if _val:
                _ph_num = _ph_key[2:]
                placeholder_text_map[f"Placeholder {_ph_num}"] = {
                    "value": str(_val), "status": "normal", "bold": False,
                }

    # Clear all annotation placeholder texts that were NOT assigned a value.
    # This prevents "Placeholder 140" etc. from showing literally in the slide.
    _ALL_ANNOTATION_KEYS = _ANNOTATION_28PT + _ANNOTATION_18PT + _ANNOTATION_PLAIN
    for _ph_key in _ALL_ANNOTATION_KEYS:
        _ph_num = _ph_key[2:]
        _ph_text = f"Placeholder {_ph_num}"
        if _ph_text not in placeholder_text_map:
            placeholder_text_map[_ph_text] = {"value": "", "status": "normal", "bold": False}

    if patient_info:
        if patient_info.get("patient_name"):
            placeholder_text_map["PatientName"] = {"value": patient_info["patient_name"], "status": "normal"}
            placeholder_text_map["Placeholder 34"] = {"value": patient_info["patient_name"], "status": "normal"}
            
        if patient_info.get("date"):
            placeholder_text_map["Date"] = {"value": patient_info["date"], "status": "normal"}
            
        if patient_info.get("complaint"):
            placeholder_text_map["Placeholder 35"] = {"value": patient_info["complaint"], "status": "normal"}
            
        if patient_info.get("gender"):
            placeholder_text_map["Placeholder 36"] = {"value": patient_info["gender"], "status": "normal"}
            
        if patient_info.get("age"):
            placeholder_text_map["Placeholder 37"] = {"value": patient_info["age"], "status": "normal"}

        # PH300: Hekim adı — "Dt. İsim Soyisim" formatında
        _doc_raw = str(patient_info.get("doctor_name", "") or "").strip()
        if _doc_raw:
            _doc_titled = " ".join(w.capitalize() for w in _doc_raw.split())
            placeholder_text_map["Placeholder 300"] = {
                "value": f"Dt. {_doc_titled}", "status": "normal", "bold": False
            }

    # Sort keys by length descending so "Placeholder 10" is processed before "Placeholder 1"
    sorted_keys = sorted(placeholder_text_map.keys(), key=len, reverse=True)
    for find_text in sorted_keys:
        info = placeholder_text_map[find_text]
        find_and_replace_text(
            prs, find_text, info["value"], info["status"],
            bold=info.get("bold", False),
            align=info.get("align"),
            font_size=info.get("font_size"),
            font_name=info.get("font_name"),
        )

    # Fix text alignment for patient info boxes on slide 0 (Slayt 1).
    # After placeholder replacement, ensure name/gender/age are left-aligned
    # and word-wrap is on so they don't overflow or appear misaligned.
    _SLIDE0_INFO_BOXES = {"TextBox 1", "TextBox 2", "TextBox 9"}
    slide0 = prs.slides[0]
    for shape in slide0.shapes:
        if shape.name in _SLIDE0_INFO_BOXES and shape.has_text_frame:
            tf = shape.text_frame
            tf.word_wrap = True
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.LEFT

    if image_paths:
        for slot_key, img_path in image_paths.items():
            if not img_path or not os.path.exists(img_path):
                continue
            mappings = IMAGE_SLOT_MAP.get(slot_key)
            if not mappings:
                print(f"  Warning: No mapping for {slot_key}")
                continue

            for mapping in mappings:
                slide_idx = mapping["slide"]
                ph_idx = mapping["placeholder_idx"]
                match_top  = mapping.get("match_top")
                match_left = mapping.get("match_left")

                if slide_idx < len(prs.slides):
                    slide = prs.slides[slide_idx]
                    success = insert_image_to_placeholder(
                        slide, ph_idx, img_path,
                        match_top=match_top, match_left=match_left,
                        fit=mapping.get('fit', 'cover'),
                    )
                    if success:
                        print(f"  [OK] Inserted {slot_key} -> Slide {slide_idx}, PH {ph_idx}")
                    else:
                        print(f"  [FAIL] Failed to insert {slot_key} -> Slide {slide_idx}, PH {ph_idx}")

    # Move the vertical blue line on Slayt 4 (slide index 3) to top=0
    # so it spans the full slide height and the photo appears fully centered
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if shape.name == "Straight Connector 9":
            shape.top = 0
            break

    # Bring text/title shapes to front on ALL slides so they are never hidden behind
    # inserted images. python-pptx adds pictures at the end of spTree (top of z-order),
    # so any existing text shapes end up behind them. Moving them to the end fixes this.
    for slide in prs.slides:
        sp_tree = slide.shapes._spTree
        text_elements = [
            shape._element for shape in slide.shapes
            if shape.has_text_frame
        ]
        for el in text_elements:
            sp_tree.remove(el)
            sp_tree.append(el)

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    blank_layout = prs.slide_layouts[6]  # blank layout

    def _insert_slide_at(prs, position):
        """Move the last-added slide to the given 0-based position."""
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        el = slides[-1]
        xml_slides.remove(el)
        xml_slides.insert(position, el)

    # ── Slayt 20 (1-tabanlı = 0-tabanlı 19): yaşa göre sil veya fotoyu yerleştir ──
    _patient_age = None
    if patient_info:
        try:
            _patient_age = float(str(patient_info.get("age", "")).strip())
        except (ValueError, TypeError):
            pass

    AFTER_COMPOSITE = 20  # 0-based index of composite slide
    if _patient_age is not None and len(prs.slides) > 19:
        if _patient_age > 18:
            # 18 yaş üstü → slayt 20'yi (0-tabanlı 19) sil
            _sldIdLst = prs.slides._sldIdLst
            _sldIdLst.remove(list(_sldIdLst)[19])
            AFTER_COMPOSITE = 19  # kompozit bir geri kaydı
            print(f"  [OK] Slayt 20 silindi (yaş {_patient_age} > 18)")
        else:
            # 18 yaş altı → sağ üst köşeye intraoral cephe fotosunu yerleştir
            _intra_path = (image_paths or {}).get("intraoral_frontal", "")
            if _intra_path and os.path.isfile(_intra_path):
                try:
                    _slide19 = prs.slides[19]
                    _left = slide_w // 2
                    _w    = slide_w // 2
                    _h    = slide_h // 2
                    _slide19.shapes.add_picture(_intra_path, _left, Emu(0), _w, _h)
                    print(f"  [OK] Slayt 20 — intraoral cephe sağ üst köşeye eklendi")
                except Exception as _e19:
                    print(f"  [WARN] Slayt 20 fotoğraf eklenemedi: {_e19}")

    next_insert = AFTER_COMPOSITE + 1  # slayt 22 pozisyonu

    # ── Slayt 22: Kapanış videosu (eğer varsa PA filminden önce eklenir) ──────
    _VIDEO_MIME_MAP = {
        '.mp4':  'video/mp4',
        '.m4v':  'video/mp4',
        '.mov':  'video/quicktime',
        '.avi':  'video/avi',
        '.wmv':  'video/x-ms-wmv',
        '.webm': 'video/webm',
        '.mkv':  'video/x-matroska',
    }
    if closing_video_path and os.path.isfile(closing_video_path):
        try:
            _ext = os.path.splitext(closing_video_path)[1].lower()
            _mime = _VIDEO_MIME_MAP.get(_ext, 'video/mp4')
            video_slide = prs.slides.add_slide(blank_layout)
            video_slide.shapes.add_movie(
                closing_video_path,
                left=0, top=0, width=slide_w, height=slide_h,
                mime_type=_mime,
            )
            _insert_slide_at(prs, next_insert)
            next_insert += 1
            print(f"  [OK] Slayt 22 — kapanış videosu eklendi")
        except Exception as ve:
            print(f"  [WARN] Video slaytı atlandı: {ve}")

    # ── Slayt 23: Posterior Anterior (PA) filmi — tam ekran ──────────────────
    pa_path = (pa_film_path or '') or (image_paths or {}).get('pa_film', '')
    if pa_path and os.path.isfile(pa_path):
        try:
            pa_slide = prs.slides.add_slide(blank_layout)
            # Fill height, preserve aspect ratio (no distortion), center horizontally
            with Image.open(pa_path) as _pa_img:
                _pa_iw, _pa_ih = _pa_img.size
            _pa_scale   = slide_h / max(_pa_ih, 1)
            _pa_pic_w   = int(_pa_iw * _pa_scale)
            _pa_pic_left = (slide_w - _pa_pic_w) // 2
            pa_slide.shapes.add_picture(pa_path, _pa_pic_left, 0, _pa_pic_w, slide_h)
            _insert_slide_at(prs, next_insert)
            print(f"  [OK] Slayt 23 — PA filmi eklendi")
        except Exception as pe:
            print(f"  [WARN] PA filmi slaytı atlandı: {pe}")

    prs.save(output_path)
    print(f"  [OK] PPTX saved to: {output_path}")
    return output_path


def inspect_template(template_path: str) -> list:
    """
    Inspect a PPTX template and return a list of all placeholders.
    Useful for mapping configuration.
    """
    prs = Presentation(template_path)
    result = []
    for i, slide in enumerate(prs.slides):
        slide_info = {
            "slide_index": i,
            "layout": slide.slide_layout.name,
            "placeholders": [],
            "shapes": [],
        }
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph = shape.placeholder_format
                slide_info["placeholders"].append({
                    "idx": ph.idx,
                    "name": shape.name,
                    "type": str(ph.type),
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                })
            else:
                shape_data = {
                    "name": shape.name,
                    "type": str(shape.shape_type),
                }
                if shape.has_text_frame:
                    shape_data["text"] = shape.text_frame.text[:100]
                slide_info["shapes"].append(shape_data)
        result.append(slide_info)
    return result
