from pptx import Presentation
p = Presentation('sunum.pptx')

print("=== SLIDE 0 Shapes ===")
for sh in p.slides[0].shapes:
    print(f"Name='{sh.name}', type={sh.shape_type}, left={sh.left}, top={sh.top}, text={sh.text_frame.text.replace(chr(10), ' ') if sh.has_text_frame else ''}")

print("\n=== All Slides with possible image slots ===")
for i in range(len(p.slides)):
    shapes = list(p.slides[i].shapes)
    if len(shapes) > 5: # likely a slide with multiple photos
        print(f"\nSlide {i} has {len(shapes)} shapes:")
        for sh in shapes:
            is_ph = "PH" if sh.is_placeholder else "Shape"
            idx = sh.placeholder_format.idx if sh.is_placeholder else "-"
            print(f"  {is_ph} name='{sh.name}' idx={idx} L={sh.left} T={sh.top} W={sh.width}")
