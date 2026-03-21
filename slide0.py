from pptx import Presentation
p = Presentation('sunum.pptx')
slide0 = p.slides[0]
for sh in slide0.shapes:
    print(f"{sh.name} - is_placeholder: {sh.is_placeholder}")
    if sh.is_placeholder:
        print(f"  idx: {sh.placeholder_format.idx}, type: {sh.placeholder_format.type}")
