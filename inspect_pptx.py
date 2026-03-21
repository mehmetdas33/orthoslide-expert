from pptx import Presentation
import sys

p = Presentation('sunum.pptx')
print("=== TEXT PLACEHOLDERS ===")
for i, sl in enumerate(p.slides):
    for sh in sl.shapes:
        if sh.has_text_frame:
            text = sh.text_frame.text.replace('\n', ' ')
            if "Placeholder" in text:
                print(f"Slide {i}: {text}")

print("\n=== IMAGES (Slides 18-23) ===")
for i in range(18, min(24, len(p.slides))):
    print(f"Slide {i}")
    for sh in p.slides[i].shapes:
        if sh.is_placeholder:
            print(f"  ph idx={sh.placeholder_format.idx} type={sh.placeholder_format.type}")
        else:
            print(f"  Shape {sh.name} type={sh.shape_type}")
sys.stdout.flush()
